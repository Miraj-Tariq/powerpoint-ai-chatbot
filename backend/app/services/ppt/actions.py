import base64
import hashlib
import os
import traceback
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List

from fastapi import UploadFile, HTTPException
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation as ppt_type
from pptx.slide import Slide
from pptx.util import Mm, Pt

from app.schemas.actions import ActionsList, ShapeParameters, ParagraphAttributes


class PPTActionsService:
    @staticmethod
    def save_ppt(presentation: UploadFile, checksum: str) -> Dict[str, str]:
        content = presentation.file.read()
        computed_checksum = hashlib.sha256(content).hexdigest()

        if checksum != computed_checksum:
            raise HTTPException(status_code=400, detail="Checksum mismatch")

        file_path = Path("current_ppt.pptx")
        with file_path.open("wb") as f:
            f.write(content)

        return {
            "message": "Presentation saved successfully",
            "filename": str(file_path)
        }


class PPTActionHandler:
    """
    Handles actions to modify a PowerPoint presentation such as
    creating/updating/deleting shapes and slides.
    """

    def __init__(
            self,
            file_name: str,
            presentation: ppt_type,
            slide: Slide,
            slide_idx: int,
            ppt_actions_GPT: ActionsList,
            selected_shape_index: Optional[int],
            attached_file: Optional[Any]) -> None:
        self.file_name = file_name
        self.presentation = presentation
        self.slide = slide
        self.slide_idx = slide_idx
        self.cwd_path = Path.cwd()
        self.ppt_actions_GPT = ppt_actions_GPT
        self.selected_shape_index = selected_shape_index
        self.attached_file = attached_file
        self.action_map = {
            "create_textbox": self._create_textbox,
            "update_textbox": self._update_textbox,
            "create_image": self._create_image,
            "update_image": self._update_image,
            "create_icon": self._create_icon,
            "update_icon": self._update_icon,
            "delete_shape": self._delete_shape
        }

    def execute_actions(self) -> None:
        try:
            for action in self.ppt_actions_GPT.actions:
                self.execute_action(
                    action_type=action.action_type.value,
                    parameters=action,
                    attached_file=self.attached_file if self.attached_file else None
                )
        except Exception as e:
            print(f"An error occurred during action execution:"
                  f"\nError Type: {type(e).__name__}"
                  f"\nError Message: {str(e)}"
                  f"\nTraceback:{traceback.format_exc()}")

    def execute_action(self,
                       action_type: str,
                       parameters: ShapeParameters,
                       attached_file: Optional[str] = None) -> None:
        """
        Executes the specified action type with provided parameters.
        """
        try:
            if action_type not in self.action_map:
                raise ValueError(f"Unsupported action type: {action_type}")

            self.action_map[action_type](parameters, attached_file)
        except Exception as e:
            print(f"An error occurred during action execution:"
                  f"\nError Type: {type(e).__name__}"
                  f"\nError Message: {str(e)}"
                  f"\nTraceback:{traceback.format_exc()}")

    def _create_textbox(self, parameters: ShapeParameters, *_) -> None:
        textbox = self.slide.shapes.add_textbox(
            Mm(parameters.left),
            Mm(parameters.top),
            Mm(parameters.width),
            Mm(parameters.height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = parameters.word_wrap if parameters.word_wrap else True
        self._add_paragraphs(text_frame, parameters.paragraphs)

    def _update_textbox(self, parameters: ShapeParameters, *_) -> None:
        shape = self.slide.shapes[self.selected_shape_index]
        shape.text_frame.clear()
        self._add_paragraphs(shape.text_frame, parameters.paragraphs)

    def _create_image(self,
                      parameters: ShapeParameters,
                      attached_file: Optional[str]) -> None:
        if not attached_file:
            raise ValueError("Attached file is required for image creation.")

        self.slide.shapes.add_picture(
            attached_file,
            Mm(parameters.left),
            Mm(parameters.top),
            Mm(parameters.width),
            Mm(parameters.height)
        )

    def _update_image(self, parameters: ShapeParameters, *_) -> None:
        shape = self.slide.shapes[self.selected_shape_index]
        if parameters.left:
            shape.left = Mm(parameters.left)
        if parameters.top:
            shape.top = Mm(parameters.top)
        if parameters.width:
            shape.width = Mm(parameters.width)
        if parameters.height:
            shape.height = Mm(parameters.height)

    def _create_icon(self, parameters: ShapeParameters, *_) -> None:
        icon_path = self.cwd_path.joinpath("resources",
                                           "icons", f"{parameters.icon_name.value}.png")
        self.slide.shapes.add_picture(
            str(icon_path),
            Mm(parameters.left),
            Mm(parameters.top),
            Mm(parameters.width),
            Mm(parameters.height)
        )

    def _update_icon(self, parameters: ShapeParameters, *_) -> None:
        self._update_image(parameters)

    def _delete_shape(self, *_) -> None:
        shape = self.slide.shapes[self.selected_shape_index]
        shape._element.getparent().remove(shape._element)

    @staticmethod
    def _add_paragraphs(text_frame, paragraphs: List[ParagraphAttributes]):
        for paragraph in paragraphs:
            p = text_frame.add_paragraph() \
                if text_frame.text else text_frame.paragraphs[0]
            p.text = paragraph.text
            font = p.font
            font.size = Pt(paragraph.font.size)
            font.bold = paragraph.font.bold
            font.italic = paragraph.font.italic
            font.underline = paragraph.font.underline

            if paragraph.font.color:
                r, g, b = tuple(int(paragraph.font.color.lstrip("#")[i:i + 2], 16)
                                for i in (0, 2, 4))
                font.color.rgb = RGBColor(r, g, b)

    def save_presentation(self, output_directory: str = "slides_ppt") -> Dict[str, str]:
        try:
            if not os.path.exists(output_directory):
                os.makedirs(output_directory)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            file_path = str(self.cwd_path
                            / output_directory
                            / f"current_ppt_{timestamp}.pptx")

            # Get the slide ID list element in the XML
            slide_id_list = self.presentation.slides._sldIdLst
            slide_indices_to_delete = [i for i in range(len(self.presentation.slides))]
            slide_indices_to_delete.remove(self.slide_idx)

            # Remove slides based on indices
            for index in sorted(slide_indices_to_delete, reverse=True):
                del slide_id_list[index]

            self.presentation.save(file_path)

            base64_encoded_ppt = self._convert_to_base64(file_path)
            return {"file_path": file_path, "base64_encoded_ppt": base64_encoded_ppt}

        except Exception as e:
            print(f"An error occurred during saving:"
                  f"\nError Type: {type(e).__name__}"
                  f"\nError Message: {str(e)}"
                  f"\nTraceback:{traceback.format_exc()}")
            return {}

    @staticmethod
    def _convert_to_base64(file_path: str) -> str:
        try:
            with open(file_path, "rb") as file:
                return base64.b64encode(file.read()).decode("utf-8")
        except Exception as e:
            print(f"Error encoding file to Base64: {str(e)}")
            return ""
