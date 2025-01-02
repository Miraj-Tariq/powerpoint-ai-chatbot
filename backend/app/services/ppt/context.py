import os
import traceback
from typing import Optional, Dict, Any

from pptx import Presentation
from pptx.presentation import Presentation as ppt_type
from pptx.slide import Slide


class SlideContext:
    def __init__(self, presentation_path: str):
        self.presentation_path = presentation_path
        self.presentation = self._load_presentation()

    def _load_presentation(self) -> ppt_type:
        try:
            return Presentation(self.presentation_path)
        except Exception as e:
            raise ValueError(f"Error loading presentation: {str(e)}")

    def get_slide(self, slide_index: int) -> Slide:
        try:
            return self.presentation.slides[slide_index]
        except IndexError:
            raise ValueError(f"Slide index {slide_index} is out of range.")


class ShapeContextExtractor:
    @staticmethod
    def extract_shape_context(slide: Slide,
                              shape_name: Optional[str] = None) -> Dict[str, Any]:
        shapes_info: list[dict] = []
        selected_shape_info = {}
        covered_areas = []

        for i, shape in enumerate(slide.shapes):
            shape_info = {
                "name": getattr(shape, "name", None),
                "shape_type": getattr(shape, "shape_type", None),
                "left": getattr(shape, "left", None),
                "top": getattr(shape, "top", None),
                "width": getattr(shape, "width", None),
                "height": getattr(shape, "height", None),
            }

            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text

            if shape_name and shape.name == shape_name:
                selected_shape_info = {
                    "actual": i,
                    "relative": len(shapes_info),
                    "info": shape_info
                }

            if all(attr is not None for attr in [shape.left,
                                                 shape.top,
                                                 shape.width,
                                                 shape.height]):
                covered_areas.append((
                    shape.top.mm,
                    shape.left.mm,
                    shape.width.mm,
                    shape.height.mm
                ))

            shapes_info.append(shape_info)

        return {
            "shapes_info": shapes_info,
            "selected_shape_info": selected_shape_info,
            "covered_areas": covered_areas
        }


class PresentationContext:
    def __init__(self, presentation_path: str):
        self.slide_context = SlideContext(presentation_path)

    def get_presentation_metadata(self) -> Dict[str, Any]:
        try:
            return {
                "file_name": os.path.basename(self.slide_context.presentation_path),
                "file_size": os.path.getsize(self.slide_context.presentation_path),
                "slide_count": len(self.slide_context.presentation.slides),
                "slide_width": self.slide_context.presentation.slide_width.mm
                if self.slide_context.presentation.slide_width else None,
                "slide_height": self.slide_context.presentation.slide_height.mm
                if self.slide_context.presentation.slide_height else None,
            }
        except Exception as e:
            raise ValueError(f"Error extracting metadata: {str(e)}")

    def get_slide_context(self,
                          slide_index: int,
                          shape_name: Optional[str] = None) -> Dict[str, Any]:
        try:
            slide = self.slide_context.get_slide(slide_index)
            shape_context = ShapeContextExtractor.extract_shape_context(
                slide,
                shape_name
            )
            return {
                "presentation_info": self.get_presentation_metadata(),
                "shapes": shape_context["shapes_info"],
                "selected_shape": shape_context["selected_shape_info"],
                "covered_areas": shape_context["covered_areas"]
            }
        except Exception as e:
            print(f"An error occurred:"
                  f"\nError Type: {type(e).__name__}"
                  f"\nError Message: {str(e)}"
                  f"\nTraceback:{traceback.format_exc()}")
            return {}
